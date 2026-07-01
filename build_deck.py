"""Generate lumen_pitch.pptx — the Lumen pitch deck.

Design system mirrors the website: warm near-black "ink" ground, gold "lumen"
light accent, teal "safe/local" accent, monospace headings for terminal DNA.
Run:  .venv/bin/python build_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ------------------------------------------------------------
INK      = RGBColor(0x14, 0x11, 0x0C)
SURFACE  = RGBColor(0x1E, 0x1A, 0x13)
SURFACE2 = RGBColor(0x27, 0x22, 0x18)
LINE     = RGBColor(0x3A, 0x32, 0x24)
LUMEN    = RGBColor(0xF4, 0xB7, 0x40)
LUMEN_HI = RGBColor(0xFF, 0xD9, 0x80)
SAFE     = RGBColor(0x4F, 0xD1, 0xC5)
FG       = RGBColor(0xF1, 0xE9, 0xD8)
MUTED    = RGBColor(0xA9, 0x9C, 0x82)
RED      = RGBColor(0xE8, 0x7A, 0x6B)

MONO = "Menlo"          # falls back gracefully off-mac
BODY = "Helvetica Neue" # falls back to Arial/Helvetica

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- helpers ------------------------------------------------------------
def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def rect(s, x, y, w, h, color, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, font, size, color, bold, tracking)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, font, size, color, bold, tracking) in para:
            r = p.add_run()
            r.text = txt
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            if tracking:
                _set_tracking(r, tracking)
    return tb


def _set_tracking(run, pts):
    """letter-spacing in points (approx). spc is in 1/100 pt."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def eyebrow(s, txt, color=LUMEN, x=Inches(0.9), y=Inches(0.62)):
    text(s, x, y, Inches(11), Inches(0.4),
         [[(txt.upper(), MONO, 12, color, True, 2.2)]])


def footer(s, n):
    text(s, Inches(0.9), Inches(7.02), Inches(4), Inches(0.3),
         [[("LUMEN", MONO, 9, MUTED, True, 2.0)]])
    text(s, Inches(11.6), Inches(7.02), Inches(0.9), Inches(0.3),
         [[(f"{n:02d}", MONO, 9, MUTED, False, 1.0)]], align=PP_ALIGN.RIGHT)


def title(s, lines, y=Inches(1.15), color=FG, size=40):
    runs = [[(ln, BODY, size, color, True, 0)] for ln in lines]
    text(s, Inches(0.9), y, Inches(11.5), Inches(2.0), runs, line_spacing=1.02, space_after=2)


def bullets(s, items, x=Inches(0.9), y=Inches(3.0), w=Inches(11.5),
            size=19, gap=14, marker=LUMEN):
    tb = s.shapes.add_textbox(x, y, w, Inches(3.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        head, sub = (it if isinstance(it, tuple) else (it, None))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        r = p.add_run(); r.text = "▸  "; r.font.name = MONO; r.font.size = Pt(size); r.font.color.rgb = marker; r.font.bold = True
        r = p.add_run(); r.text = head; r.font.name = BODY; r.font.size = Pt(size); r.font.color.rgb = FG; r.font.bold = True
        if sub:
            r = p.add_run(); r.text = "   " + sub; r.font.name = BODY; r.font.size = Pt(size); r.font.color.rgb = MUTED; r.font.bold = False
    return tb


def glow(s, cx, cy, r, color):
    """soft luminous disc behind hero elements."""
    sp = rect(s, cx - r, cy - r, r * 2, r * 2, color, shape=MSO_SHAPE.OVAL)
    _soft(sp, color)
    return sp


def _soft(sp, color):
    # low-opacity fill via alpha on solidFill
    solidFill = sp.fill._xPr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    alpha = srgb.makeelement(qn('a:alpha'), {'val': '14000'})
    srgb.append(alpha)


N = 0
def num():
    global N
    N += 1
    return N


# ========================================================================
# 1 — TITLE
# ========================================================================
s = slide()
glow(s, Inches(10.6), Inches(2.2), Inches(2.6), LUMEN)
rect(s, Inches(0.9), Inches(2.35), Inches(0.55), Inches(0.09), LUMEN)
text(s, Inches(0.9), Inches(2.55), Inches(11.5), Inches(2.2), [
    [("Lumen", BODY, 76, FG, True, 0)],
], line_spacing=1.0)
text(s, Inches(0.92), Inches(3.95), Inches(11.5), Inches(1.0), [
    [("A privacy-first coding harness.", BODY, 26, LUMEN_HI, False, 0)],
])
text(s, Inches(0.92), Inches(4.75), Inches(10.8), Inches(1.2), [
    [("Run local LLMs offline, or bring your own OpenRouter key for the whole", BODY, 17, MUTED, False, 0)],
    [("catalogue. No daily limits. Your code never leaves your machine.", BODY, 17, MUTED, False, 0)],
], line_spacing=1.2, space_after=2)
text(s, Inches(0.9), Inches(0.62), Inches(11), Inches(0.4),
     [[("HACKMEE  ·  DEVELOPER TOOLING", MONO, 12, LUMEN, True, 2.4)]])
footer(s, num())

# ========================================================================
# 2 — THE PROBLEM
# ========================================================================
s = slide()
eyebrow(s, "The problem")
title(s, ["“You’ve hit your daily limit.”"])
text(s, Inches(0.9), Inches(2.0), Inches(11.4), Inches(0.9),
     [[("The moment you get into flow, the hosted coding tools tap out — and", BODY, 18, MUTED, False, 0)],
      [("that’s not the only tax you’re paying.", BODY, 18, MUTED, False, 0)]], line_spacing=1.2, space_after=2)

cards = [
    ("⏳", "Daily limits", "Claude Code and Codex lock you out mid-task. Your productivity is capped by someone else’s quota.", RED),
    ("🔒", "Your code leaves", "Every file, prompt and secret is uploaded to a vendor’s cloud. For proprietary code, that’s a non-starter.", LUMEN),
    ("⛓", "Vendor lock-in", "One provider, one model, one price. No fallback when it’s down, slow, or discontinued.", SAFE),
]
cw = Inches(3.72); gap = Inches(0.28); x0 = Inches(0.9); cy = Inches(3.15); ch = Inches(2.9)
for i, (icon, head, body, accent) in enumerate(cards):
    x = x0 + i * (cw + gap)
    rect(s, x, cy, cw, ch, SURFACE, line=LINE, line_w=1)
    rect(s, x, cy, cw, Inches(0.09), accent)
    text(s, x + Inches(0.3), cy + Inches(0.34), cw - Inches(0.6), Inches(0.6), [[(icon, BODY, 26, accent, False, 0)]])
    text(s, x + Inches(0.3), cy + Inches(1.0), cw - Inches(0.6), Inches(0.5), [[(head, BODY, 20, FG, True, 0)]])
    text(s, x + Inches(0.3), cy + Inches(1.55), cw - Inches(0.6), Inches(1.2), [[(body, BODY, 14, MUTED, False, 0)]], line_spacing=1.18)
footer(s, num())

# ========================================================================
# 3 — THE IDEA
# ========================================================================
s = slide()
glow(s, Inches(6.6), Inches(3.6), Inches(3.2), LUMEN)
eyebrow(s, "The idea")
text(s, Inches(1.3), Inches(2.5), Inches(10.7), Inches(2.6), [
    [("What if your coding agent ran on ", BODY, 40, FG, True, 0), ("your", BODY, 40, LUMEN, True, 0), ("", BODY, 40, FG, True, 0)],
    [("machine — and answered only to you?", BODY, 40, FG, True, 0)],
], align=PP_ALIGN.CENTER, line_spacing=1.08, space_after=2)
text(s, Inches(1.3), Inches(4.7), Inches(10.7), Inches(0.8),
     [[("Same agentic power. None of the limits, none of the exposure.", BODY, 19, MUTED, False, 0)]],
     align=PP_ALIGN.CENTER)
footer(s, num())

# ========================================================================
# 4 — MEET LUMEN / TWO MODES
# ========================================================================
s = slide()
eyebrow(s, "The solution")
title(s, ["One harness. Two modes."])
text(s, Inches(0.9), Inches(2.0), Inches(11.4), Inches(0.6),
     [[("A real terminal coding agent you own — it reads, edits, runs and searches your code.", BODY, 18, MUTED, False, 0)]])

# two big mode cards
mw = Inches(5.66); mh = Inches(3.15); my = Inches(3.05)
# local
rect(s, Inches(0.9), my, mw, mh, SURFACE, line=LINE)
rect(s, Inches(0.9), my, Inches(0.09), mh, SAFE)
text(s, Inches(1.25), my + Inches(0.3), mw - Inches(0.7), Inches(0.4), [[("LOCAL · OFFLINE", MONO, 13, SAFE, True, 2.0)]])
text(s, Inches(1.25), my + Inches(0.78), mw - Inches(0.7), Inches(0.5), [[("Private by default", BODY, 24, FG, True, 0)]])
for i, t in enumerate(["Ollama · LM Studio · llama.cpp · vLLM", "Code never touches the network — at all", "$0 per token, unlimited runs", "Works on a plane, in a vault, anywhere"]):
    text(s, Inches(1.25), my + Inches(1.35) + i*Inches(0.42), mw - Inches(0.7), Inches(0.4),
         [[("•  ", BODY, 15, SAFE, True, 0), (t, BODY, 15, FG if i>0 else FG, False, 0)]])
# openrouter
x2 = Inches(0.9) + mw + Inches(0.29)
rect(s, x2, my, mw, mh, SURFACE, line=LINE)
rect(s, x2, my, Inches(0.09), mh, LUMEN)
text(s, x2 + Inches(0.35), my + Inches(0.3), mw - Inches(0.7), Inches(0.4), [[("OPENROUTER · YOUR KEY", MONO, 13, LUMEN, True, 2.0)]])
text(s, x2 + Inches(0.35), my + Inches(0.78), mw - Inches(0.7), Inches(0.5), [[("The whole catalogue", BODY, 24, FG, True, 0)]])
for i, t in enumerate(["Claude, GPT, Gemini, Llama, Qwen, DeepSeek…", "No daily limits — pay only per token used", "Switch models without leaving the REPL", "One key you control, transparent cost"]):
    text(s, x2 + Inches(0.35), my + Inches(1.35) + i*Inches(0.42), mw - Inches(0.7), Inches(0.4),
         [[("•  ", BODY, 15, LUMEN, True, 0), (t, BODY, 15, FG, False, 0)]])
footer(s, num())

# ========================================================================
# 5 — HOW IT WORKS
# ========================================================================
s = slide()
eyebrow(s, "How it works")
title(s, ["A real agent loop — not a chat box."])
# loop chips
steps = [("prompt", MUTED), ("stream reply", FG), ("call tools", LUMEN), ("run + observe", SAFE), ("repeat", FG)]
x = Inches(0.9); yv = Inches(2.35)
for i, (t, c) in enumerate(steps):
    w = Inches(2.0)
    rect(s, x, yv, w, Inches(0.62), SURFACE, line=LINE)
    text(s, x, yv + Inches(0.16), w, Inches(0.4), [[(t, MONO, 13, c, True, 0.5)]], align=PP_ALIGN.CENTER)
    if i < len(steps)-1:
        text(s, x + w, yv + Inches(0.12), Inches(0.42), Inches(0.4), [[("→", BODY, 20, LUMEN, True, 0)]], align=PP_ALIGN.CENTER)
    x += w + Inches(0.42)
text(s, Inches(0.9), Inches(3.25), Inches(11.4), Inches(0.5),
     [[("Loops until the model stops calling tools. Every write or shell command asks permission first.", BODY, 16, MUTED, False, 0)]])
# tools grid
tools = [("read_file","read"),("write_file","write ⚠"),("edit_file","edit ⚠"),("list_dir","list"),
         ("run_bash","shell ⚠"),("search","grep"),("find_files","glob")]
tw = Inches(1.5); th = Inches(1.0); gx = Inches(0.15); tx0 = Inches(0.9); ty = Inches(4.0)
for i, (name, kind) in enumerate(tools):
    x = tx0 + i*(tw+gx)
    rect(s, x, ty, tw, th, SURFACE2, line=LINE)
    warn = "⚠" in kind
    text(s, x, ty + Inches(0.2), tw, Inches(0.35), [[(name, MONO, 12, LUMEN if not warn else LUMEN_HI, True, 0)]], align=PP_ALIGN.CENTER)
    text(s, x, ty + Inches(0.55), tw, Inches(0.3), [[(kind, MONO, 11, MUTED, False, 0)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.9), Inches(5.2), Inches(11), Inches(0.4),
     [[("⚠ = asks approval before it runs.  Toggle /auto or --yolo to run unattended.", MONO, 12, MUTED, False, 0.5)]])
footer(s, num())

# ========================================================================
# 6 — FEATURES OVERVIEW
# ========================================================================
s = slide()
eyebrow(s, "What's inside")
title(s, ["Six things hosted tools won’t do."], size=32)
feats = [
    ("🛡", "Secret Guard", "Blocks API keys & .env values before they reach the cloud. Skipped offline.", LUMEN),
    ("📝", "Diff before write", "A colored diff of every change — you approve before it happens.", SAFE),
    ("↩", "Instant undo", "/undo reverts the last change, whether it edited or created a file.", LUMEN),
    ("💾", "Local sessions", "--continue / --resume. Your history lives on your disk, not a server.", SAFE),
    ("🧠", "Project memory", "Auto-loads LUMEN.md for your conventions; /init writes one for you.", LUMEN),
    ("🔁", "Auto-fallback", "Model down? Fail over to another provider — with an honest warning.", SAFE),
]
fcw, fch, fgx, fgy = Inches(3.72), Inches(1.78), Inches(0.28), Inches(0.22)
fx0, fy0 = Inches(0.9), Inches(2.35)
for i, (icon, head, desc, accent) in enumerate(feats):
    col, row = i % 3, i // 3
    x = fx0 + col * (fcw + fgx)
    y = fy0 + row * (fch + fgy)
    rect(s, x, y, fcw, fch, SURFACE, line=LINE)
    rect(s, x, y, fcw, Inches(0.08), accent)
    text(s, x + Inches(0.28), y + Inches(0.26), fcw - Inches(0.5), Inches(0.4), [[(icon + "  " + head, BODY, 17, FG, True, 0)]])
    text(s, x + Inches(0.28), y + Inches(0.78), fcw - Inches(0.55), Inches(0.9), [[(desc, BODY, 12.5, MUTED, False, 0)]], line_spacing=1.15)
footer(s, num())

# ========================================================================
# 7 — SECRET GUARD SPOTLIGHT
# ========================================================================
s = slide()
glow(s, Inches(11.2), Inches(1.8), Inches(2.4), RED)
eyebrow(s, "Flagship feature · privacy", color=RED)
title(s, ["Your secrets never leave by accident."], size=30)
text(s, Inches(0.9), Inches(1.95), Inches(11.4), Inches(0.7),
     [[("Before Lumen sends anything to a cloud model, it scans every message for keys, ", BODY, 16, MUTED, False, 0)],
      [("private keys and .env values — and lets you block or redact. Offline, it never even runs.", BODY, 16, MUTED, False, 0)]],
     line_spacing=1.2, space_after=2)
gx, gy, gw, gh = Inches(0.9), Inches(3.0), Inches(11.5), Inches(3.55)
rect(s, gx, gy, gw, gh, RGBColor(0x0D, 0x0B, 0x08), line=LINE)
rect(s, gx, gy, gw, Inches(0.42), SURFACE2)
for i, c in enumerate([RED, LUMEN, SAFE]):
    rect(s, gx + Inches(0.28) + i * Inches(0.28), gy + Inches(0.15), Inches(0.13), Inches(0.13), c, shape=MSO_SHAPE.OVAL)
text(s, gx, gy + Inches(0.11), gw, Inches(0.3), [[("lumen — claude-sonnet-4.5 · openrouter · cloud", MONO, 11, MUTED, False, 0.5)]], align=PP_ALIGN.CENTER)
glines = [
    ("› read .env and tell me what's configured", FG, True),
    ("  ⚙ read_file   .env", SAFE, False),
    ("  🛡 Secret Guard — about to send 2 secrets to openrouter (cloud)", RED, True),
    ("     • OpenRouter API key   sk-…ba (45 chars)   [in a file result]", MUTED, False),
    ("     • AWS access key id    AKI…LE (20 chars)   [in a file result]", MUTED, False),
    ("  redact & send (r) / send anyway (s) / cancel (c):  c", LUMEN, False),
    ("  ✗ Cancelled — nothing was sent.", RED, True),
    ("  (in local mode this check never runs — your code never leaves)", MUTED, False),
]
tb = s.shapes.add_textbox(gx + Inches(0.4), gy + Inches(0.62), gw - Inches(0.8), gh - Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True
for i, (t, c, b) in enumerate(glines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6); p.line_spacing = 1.05
    r = p.add_run(); r.text = t; r.font.name = MONO; r.font.size = Pt(13.5); r.font.color.rgb = c; r.font.bold = b
# airgap "goes further" callout
rect(s, gx, Inches(6.62), gw, Inches(0.34), SURFACE2, line=SAFE)
text(s, gx + Inches(0.2), Inches(6.66), gw - Inches(0.4), Inches(0.3),
     [[("🔒 Airgap mode goes further — one flag blocks every socket. Local keeps working; nothing can leave.", MONO, 12, SAFE, True, 0.3)]])
footer(s, num())

# ========================================================================
# 8 — PROOF (real transcript)
# ========================================================================
s = slide()
eyebrow(s, "Proof · running on a local model, offline")
title(s, ["It ships code, recovers from errors, verifies itself."], size=30)
# terminal panel
tx, ty, tw2, th2 = Inches(0.9), Inches(2.25), Inches(11.5), Inches(4.35)
rect(s, tx, ty, tw2, th2, RGBColor(0x0D,0x0B,0x08), line=LINE)
rect(s, tx, ty, tw2, Inches(0.42), SURFACE2)
for i, c in enumerate([RED, LUMEN, SAFE]):
    rect(s, tx + Inches(0.28) + i*Inches(0.28), ty + Inches(0.15), Inches(0.13), Inches(0.13), c, shape=MSO_SHAPE.OVAL)
text(s, tx, ty + Inches(0.11), tw2, Inches(0.3), [[("lumen — qwen3-coder:30b  ·  local  ·  offline", MONO, 11, MUTED, False, 0.5)]], align=PP_ALIGN.CENTER)
lines = [
    ("› edit greet.py: add farewell(name) and print it, then run the file", FG, True),
    ("  ⚙ read_file {\"path\": \"greet.py\"}", SAFE, False),
    ("  ⚙ edit_file {\"path\": \"greet.py\", ...}   → 1 occurrence replaced", SAFE, False),
    ("  ⚙ run_bash {\"command\": \"python greet.py\"}", SAFE, False),
    ("    /bin/sh: python: command not found   [exit 127]", RED, False),
    ("  ⚙ run_bash {\"command\": \"python3 greet.py\"}   ← recovered on its own", LUMEN, False),
    ("    Hello, world from Lumen!", FG, False),
    ("    Goodbye, world!            [exit 0]", FG, False),
    ("  lumen  Done — added farewell(), updated __main__, verified output.", LUMEN_HI, False),
]
tb = s.shapes.add_textbox(tx + Inches(0.4), ty + Inches(0.62), tw2 - Inches(0.8), th2 - Inches(0.9))
tf = tb.text_frame; tf.word_wrap = True
for i, (t, c, b) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(5); p.line_spacing = 1.05
    r = p.add_run(); r.text = t; r.font.name = MONO; r.font.size = Pt(13.5); r.font.color.rgb = c; r.font.bold = b
footer(s, num())

# ========================================================================
# 7 — THE NUMBERS (cost / limits)
# ========================================================================
s = slide()
eyebrow(s, "The numbers · cost & limits")
title(s, ["Same model, same tokens. Different bill."], size=32)
text(s, Inches(0.9), Inches(1.95), Inches(11.4), Inches(0.7),
     [[("On the same task, the same model reads the same tokens. What changes is ", BODY, 16, MUTED, False, 0),
       ("who bills you", BODY, 16, LUMEN, True, 0),
       (" and ", BODY, 16, MUTED, False, 0),
       ("whether you get locked out", BODY, 16, LUMEN, True, 0),
       (".", BODY, 16, MUTED, False, 0)]], line_spacing=1.2)

# comparison bars: "a heavy day of coding"
def cost_row(y, label, val_txt, frac, color, note):
    text(s, Inches(0.9), y, Inches(3.1), Inches(0.4), [[(label, BODY, 15, FG, True, 0)]])
    bx = Inches(4.1); bw = Inches(6.4)
    rect(s, bx, y + Inches(0.03), bw, Inches(0.34), SURFACE, line=LINE)
    if frac > 0:
        rect(s, bx, y + Inches(0.03), Emu(int(bw * frac)), Inches(0.34), color)
    text(s, bx + bw + Inches(0.2), y, Inches(2.5), Inches(0.4), [[(val_txt, MONO, 14, color, True, 0)]])
    text(s, Inches(0.9), y + Inches(0.4), Inches(3.1), Inches(0.3), [[(note, MONO, 11, MUTED, False, 0)]])

text(s, Inches(0.9), Inches(2.95), Inches(11), Inches(0.35), [[("BILLABLE $ OVER A HEAVY DAY OF CODING", MONO, 12, MUTED, True, 1.6)]])
cost_row(Inches(3.5), "Hosted subscription", "capped → locked out", 0.62, RED, "hit the daily quota, wait")
cost_row(Inches(4.4), "Lumen · OpenRouter", "pay-per-token, no cap", 0.4, LUMEN, "your key, transparent price")
cost_row(Inches(5.3), "Lumen · local", "$0 · unlimited", 1.0, SAFE, "runs entirely on your box")
text(s, Inches(0.9), Inches(6.35), Inches(11.4), Inches(0.5),
     [[("Illustrative — exact spend depends on model & usage. The facts: local bills $0 with no cap; hosted plans throttle.", MONO, 11, MUTED, False, 0)]])
footer(s, num())

# ========================================================================
# 8 — LEAN CONTEXT (measured)
# ========================================================================
s = slide()
eyebrow(s, "The numbers · lean context (measured)")
title(s, ["~1,088 tokens of overhead. Measured, not guessed."], size=30)
text(s, Inches(0.9), Inches(1.95), Inches(11.4), Inches(0.7),
     [[("Every turn, a harness prepends a system prompt + tool definitions before your code. ", BODY, 16, MUTED, False, 0),
       ("Lumen keeps that tiny", BODY, 16, LUMEN, True, 0),
       (" — more of your context window and budget goes to the actual work.", BODY, 16, MUTED, False, 0)]], line_spacing=1.2)

# big stat + breakdown
rect(s, Inches(0.9), Inches(3.0), Inches(3.5), Inches(3.0), SURFACE, line=LINE)
rect(s, Inches(0.9), Inches(3.0), Inches(3.5), Inches(0.09), LUMEN)
text(s, Inches(0.9), Inches(3.45), Inches(3.5), Inches(1.1), [[("1,088", MONO, 58, LUMEN_HI, True, 0)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.9), Inches(4.7), Inches(3.5), Inches(0.4), [[("tokens / turn", MONO, 15, FG, True, 1.0)]], align=PP_ALIGN.CENTER)
text(s, Inches(0.9), Inches(5.15), Inches(3.5), Inches(0.6), [[("296 system prompt", MONO, 12, MUTED, False, 0)],[("792 · 7 tool schemas", MONO, 12, MUTED, False, 0)]], align=PP_ALIGN.CENTER, space_after=2)

# per-tool bars
tools_tok = [("read_file",132),("edit_file",155),("run_bash",113),("search",110),("find_files",101),("write_file",98),("list_dir",75)]
bx = Inches(4.9); by = Inches(3.1); maxw = Inches(5.2); mx = 155
text(s, bx, by - Inches(0.05), Inches(6), Inches(0.35), [[("TOKENS PER TOOL SCHEMA", MONO, 12, MUTED, True, 1.6)]])
for i, (name, tok) in enumerate(tools_tok):
    y = by + Inches(0.45) + i*Inches(0.4)
    text(s, bx, y, Inches(1.5), Inches(0.3), [[(name, MONO, 12, FG, False, 0)]])
    rect(s, bx + Inches(1.5), y + Inches(0.02), Emu(int(maxw * tok/mx)), Inches(0.22), LUMEN if name not in ("write_file","run_bash","edit_file") else LUMEN_HI)
    text(s, bx + Inches(1.5) + Emu(int(maxw * tok/mx)) + Inches(0.1), y - Inches(0.02), Inches(1), Inches(0.3), [[(str(tok), MONO, 11, MUTED, False, 0)]])
footer(s, num())

# ========================================================================
# 9 — COMPARISON TABLE
# ========================================================================
s = slide()
eyebrow(s, "At a glance")
title(s, ["Lumen vs. hosted coding tools"], size=32)

rows = [
    ("", "Hosted (Claude Code / Codex)", "Lumen"),
    ("Daily usage limits", "Yes — locked out at the cap", "None (local) · none beyond your balance"),
    ("Where your code goes", "Uploaded to a vendor cloud", "Stays on your machine (local mode)"),
    ("Provable offline", "No", "Airgap blocks all egress at the socket layer"),
    ("Secret leak protection", "None — it’s all uploaded", "Guard blocks keys before they leave"),
    ("See changes first", "Hidden / varies", "Colored diff + approve, then /undo"),
    ("Cost model", "Fixed subscription, capped", "Free (local) or pay-per-token"),
    ("Model choice", "One vendor’s models", "Whole OpenRouter catalogue + any local"),
    ("If the model is down", "You wait", "Auto-fallback to another provider"),
    ("Works offline", "No", "Yes"),
    ("Context overhead / turn", "Large, hidden", "~1,088 tokens, measured"),
    ("You control it", "No", "Open, self-hosted, yours"),
]
ty = Inches(2.05); rh = Inches(0.40); c0 = Inches(0.9); c1w = Inches(3.5); c2w = Inches(4.3); c3w = Inches(4.6)
for i, (a, b, c) in enumerate(rows):
    y = ty + i*rh
    header = (i == 0)
    if header:
        rect(s, c0, y, c1w+c2w+c3w, rh, SURFACE2)
    elif i % 2 == 0:
        rect(s, c0, y, c1w+c2w+c3w, rh, SURFACE)
    text(s, c0 + Inches(0.15), y + Inches(0.13), c1w-Inches(0.2), rh, [[(a, MONO, 12.5, MUTED if not header else FG, header, 0.5)]])
    text(s, c0 + c1w + Inches(0.15), y + Inches(0.13), c2w-Inches(0.2), rh, [[(b, BODY, 13, (MUTED if not header else RED), header, 0)]])
    text(s, c0 + c1w + c2w + Inches(0.15), y + Inches(0.13), c3w-Inches(0.2), rh, [[(c, BODY, 13, (FG if not header else LUMEN), header, 0)]])
    if not header:
        rect(s, c0, y, c1w+c2w+c3w, Emu(9525), LINE)  # top hairline
footer(s, num())

# ========================================================================
# 10 — ARCHITECTURE
# ========================================================================
s = slide()
eyebrow(s, "Under the hood")
title(s, ["One client for local and cloud."], size=32)
text(s, Inches(0.9), Inches(1.95), Inches(11.4), Inches(0.7),
     [[("Ollama, LM Studio, llama.cpp and OpenRouter all speak the OpenAI ", BODY, 16, MUTED, False, 0),
       ("/chat/completions", MONO, 15, SAFE, False, 0),
       (" dialect. So a new backend is a config entry — not new code.", BODY, 16, MUTED, False, 0)]], line_spacing=1.2)
arch = [
    ("cli.py", "argument parsing · REPL · slash commands"),
    ("agent.py", "the loop: stream → run tools → repeat"),
    ("providers/openai_compat.py", "one streaming client · tool-call assembly"),
    ("tools/  (fs · shell · search)", "7 tools with JSON schemas + a permission gate"),
    ("ui/console.py", "rich streaming, tool display, approvals"),
    ("config.py", "providers, models, keys — ~/.lumen/config.json"),
]
ay = Inches(2.95)
for i, (mod, desc) in enumerate(arch):
    y = ay + i*Inches(0.62)
    rect(s, Inches(0.9), y, Inches(4.7), Inches(0.5), SURFACE, line=LINE)
    text(s, Inches(1.1), y + Inches(0.11), Inches(4.4), Inches(0.35), [[(mod, MONO, 13.5, LUMEN, True, 0)]])
    text(s, Inches(5.85), y + Inches(0.12), Inches(6.5), Inches(0.35), [[(desc, BODY, 14, MUTED, False, 0)]])
footer(s, num())

# ========================================================================
# 11 — CTA / GET STARTED
# ========================================================================
s = slide()
glow(s, Inches(11.0), Inches(5.6), Inches(2.6), LUMEN)
eyebrow(s, "Get started")
text(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(1.4), [
    [("Your machine.", BODY, 46, FG, True, 0)],
    [("Your models. Your rules.", BODY, 46, LUMEN, True, 0)],
], line_spacing=1.05, space_after=2)

# code block
cx, cyy, cw2, ch2 = Inches(0.9), Inches(3.6), Inches(7.6), Inches(2.7)
rect(s, cx, cyy, cw2, ch2, RGBColor(0x0D,0x0B,0x08), line=LINE)
cmds = [
    ("# local & offline — the default", MUTED),
    ("$ ollama pull qwen3-coder:30b", SAFE),
    ("$ lumen", FG),
    ("", FG),
    ("# cloud — your own key, full catalogue", MUTED),
    ("$ export OPENROUTER_API_KEY=sk-or-...", LUMEN),
    ("$ lumen -p openrouter", FG),
]
tb = s.shapes.add_textbox(cx + Inches(0.35), cyy + Inches(0.3), cw2 - Inches(0.6), ch2 - Inches(0.5))
tf = tb.text_frame; tf.word_wrap = True
for i, (t, c) in enumerate(cmds):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(6); p.line_spacing = 1.05
    r = p.add_run(); r.text = t if t else " "; r.font.name = MONO; r.font.size = Pt(15); r.font.color.rgb = c
# side points
sx = Inches(8.9)
for i, t in enumerate(["No daily limits", "Code stays local", "Any model, any time", "Open & self-hosted"]):
    y = Inches(3.7) + i*Inches(0.62)
    text(s, sx, y, Inches(0.5), Inches(0.4), [[("◆", BODY, 15, LUMEN, True, 0)]])
    text(s, sx + Inches(0.4), y, Inches(3.6), Inches(0.4), [[(t, BODY, 18, FG, True, 0)]])
footer(s, num())

prs.save("lumen_pitch.pptx")
print(f"Saved lumen_pitch.pptx — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
